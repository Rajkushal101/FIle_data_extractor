"""
PPTX (PowerPoint) Extraction Service
"""
import logging
from typing import Dict
from pptx import Presentation
import base64
from app.services.math_detector import math_detector

logger = logging.getLogger(__name__)


class PPTXExtractor:
    """Extract text and images from PowerPoint presentations"""
    
    async def extract(self, file_path: str) -> Dict:
        """
        Extract content from PPTX file
        
        Args:
            file_path: Path to PPTX file
            
        Returns:
            Dict with text, images, and metadata
        """
        try:
            prs = Presentation(file_path)
            
            slides_content = []
            all_images = []
            
            for slide_num, slide in enumerate(prs.slides):
                slide_text = []
                
                # Extract text from shapes
                for shape in slide.shapes:
                    shape_text = getattr(shape, "text", "")
                    if isinstance(shape_text, str) and shape_text.strip():
                        slide_text.append(shape_text)
                    
                    # Extract images
                    if shape.shape_type == 13:  # Picture type
                        try:
                            image = getattr(shape, "image", None)
                            if image is None:
                                continue
                            img_bytes = image.blob
                            img_base64 = base64.b64encode(img_bytes).decode('utf-8')
                            all_images.append({
                                "slide": slide_num + 1,
                                "data": img_base64
                            })
                        except Exception:
                            pass
                
                if slide_text:
                    slides_content.append(f"### Slide {slide_num + 1} ###\n" + "\n".join(slide_text))
            
            full_text = "\n\n".join(slides_content)

            return {
                "text": full_text,
                "images": all_images,
                "math_expressions": math_detector.extract_math_expressions(full_text),
                "math_text": math_detector.to_math_text(full_text),
                "metadata": {
                    "slide_count": len(prs.slides),
                    "image_count": len(all_images)
                }
            }
            
        except Exception as e:
            logger.error(f"Error extracting PPTX: {str(e)}")
            raise